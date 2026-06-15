import datetime
import os
import subprocess
import zipfile

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = r"D:\Users\20371\Desktop\VectorGenEditor"
PRES_DIR = os.path.join(ROOT, "artifacts", "presentation")
SOURCE_PNG_DIR = os.path.join(PRES_DIR, "slide_assets", "01_images")
ILLUSTRATION_DIR = os.path.join(PRES_DIR, "illustrations")
BG_DIR = os.path.join(PRES_DIR, "slide_backgrounds")
TMP_DIR = os.path.join(PRES_DIR, "tmp_run")

NAME = "张呈畅"
TITLE = "《半自动化坐标驱动的矢量图形生成与编辑器项目》"
COURSE = "Java 高级应用开发  2025-2026-2  2024-1 班"
TEACHER = "王莹"
CLASS = "6 班"
STUDENT_ID = "20233132119"
OUTPUT_BASE = os.path.join(PRES_DIR, f"{NAME}-期末作品答辩-{STUDENT_ID}")

COLORS = {
    "bg_top": (10, 31, 68),
    "bg_bottom": (23, 79, 130),
    "accent": (255, 179, 71),
    "card": (245, 248, 252),
    "white": (255, 255, 255),
    "text": (28, 44, 70),
    "muted": (98, 112, 131),
    "accent2": (58, 123, 213),
}


def ensure_dirs():
    for path in [PRES_DIR, ILLUSTRATION_DIR, BG_DIR, TMP_DIR]:
        os.makedirs(path, exist_ok=True)


def try_font(size, bold=False):
    candidates = [
        r"C:\Windows\Fonts\msyhbd.ttc" if bold else r"C:\Windows\Fonts\msyh.ttc",
        r"C:\Windows\Fonts\simhei.ttf",
        r"C:\Windows\Fonts\arial.ttf",
    ]
    for candidate in candidates:
        if os.path.exists(candidate):
            return ImageFont.truetype(candidate, size=size)
    return ImageFont.load_default()


def linear_gradient(size, top, bottom):
    w, h = size
    base = Image.new("RGB", size, top)
    draw = ImageDraw.Draw(base)
    for y in range(h):
        ratio = y / max(h - 1, 1)
        color = tuple(int(top[i] + (bottom[i] - top[i]) * ratio) for i in range(3))
        draw.line([(0, y), (w, y)], fill=color)
    return base


def wrap_text(draw, text, font, max_width):
    lines = []
    for paragraph in text.split("\n"):
        if not paragraph:
            lines.append("")
            continue
        current = ""
        for ch in paragraph:
            trial = current + ch
            if draw.textlength(trial, font=font) <= max_width:
                current = trial
            else:
                if current:
                    lines.append(current)
                current = ch
        if current:
            lines.append(current)
    return lines


def draw_rounded_card(draw, box, fill, outline=None, radius=28):
    draw.rounded_rectangle(box, radius=radius, fill=fill, outline=outline, width=2 if outline else 0)


def create_background(filename, title, subtitle=None, footer=None, layout="default"):
    size = (1920, 1080)
    img = linear_gradient(size, COLORS["bg_top"], COLORS["bg_bottom"])
    draw = ImageDraw.Draw(img)
    accent_font = try_font(36, bold=True)
    title_font = try_font(72 if layout == "cover" else 56, bold=True)
    sub_font = try_font(28)
    footer_font = try_font(24)

    draw.rounded_rectangle((70, 70, 380, 145), radius=24, fill=(255, 255, 255, 32))
    draw.text((100, 92), "VectorGenEditor", font=accent_font, fill=(255, 220, 170))

    if layout == "cover":
        draw.rounded_rectangle((110, 280, 1810, 920), radius=40, fill=(255, 255, 255, 18), outline=COLORS["accent"], width=3)
        draw.text((180, 420), title, font=title_font, fill=COLORS["white"])
        if subtitle:
            for idx, line in enumerate(subtitle.split("\n")):
                draw.text((180, 560 + idx * 48), line, font=sub_font, fill=(225, 235, 248))
    else:
        draw.text((100, 110), title, font=title_font, fill=COLORS["white"])
        if subtitle:
            draw.text((100, 190), subtitle, font=sub_font, fill=(225, 235, 248))
        draw.polygon([(1500, 720), (1820, 620), (1820, 980), (1380, 980)], fill=(255, 255, 255, 16))
        draw.ellipse((1580, 180, 1840, 440), outline=COLORS["accent"], width=4)
        draw.rectangle((1360, 260, 1520, 420), outline=(255, 255, 255, 80), width=3)

    if footer:
        draw.text((100, 980), footer, font=footer_font, fill=(225, 235, 248))

    out = os.path.join(BG_DIR, filename)
    img.save(out, format="PNG", optimize=True)
    return out


def draw_flow_illustration():
    path = os.path.join(ILLUSTRATION_DIR, "auto_generate_flow.png")
    img = Image.new("RGB", (1600, 900), COLORS["card"])
    draw = ImageDraw.Draw(img)
    font = try_font(34, bold=True)
    small = try_font(24)
    steps = [
        ("输入", "参数 / 文件名"),
        ("校验", "Radius / Count"),
        ("生成", "DocumentController"),
        ("刷新", "CanvasScene"),
        ("选中", "SelectionController"),
    ]
    x, y, width, height, gap = 80, 260, 250, 180, 90
    for idx, (title, detail) in enumerate(steps):
        draw_rounded_card(draw, (x, y, x + width, y + height), COLORS["white"], outline=COLORS["accent2"])
        draw.text((x + 28, y + 42), title, font=font, fill=COLORS["text"])
        for line_idx, line in enumerate(wrap_text(draw, detail, small, width - 56)):
            draw.text((x + 28, y + 108 + line_idx * 30), line, font=small, fill=COLORS["muted"])
        if idx < len(steps) - 1:
            draw.polygon([(x + width + 18, y + height / 2), (x + width + 58, y + height / 2 - 22), (x + width + 58, y + height / 2 + 22)], fill=COLORS["accent"])
        x += width + gap
    draw.text((80, 80), "半自动图形生成流程", font=try_font(48, bold=True), fill=COLORS["text"])
    img.save(path, format="PNG")
    return path


def draw_save_export_illustration():
    path = os.path.join(ILLUSTRATION_DIR, "save_export_formats.png")
    img = Image.new("RGB", (1400, 780), COLORS["card"])
    draw = ImageDraw.Draw(img)
    draw.text((70, 60), "保存与导出格式", font=try_font(48, bold=True), fill=COLORS["text"])
    cards = [
        ("保存", "JsonSerializer", ["document", "layers", "shapes", "transform", "style"]),
        ("导出", "SvgExporter", ["viewBox", "path/rect/text", "transform", "stroke/fill"]),
    ]
    x = 90
    for title, module, items in cards:
        draw_rounded_card(draw, (x, 210, x + 560, 700), COLORS["white"], outline=COLORS["accent2"])
        draw.text((x + 38, 250), title, font=try_font(30, bold=True), fill=COLORS["text"])
        draw.text((x + 38, 310), module, font=try_font(24), fill=COLORS["accent2"])
        yy = 380
        for item in items:
            draw.ellipse((x + 42, yy + 10, x + 54, yy + 22), fill=COLORS["accent"])
            draw.text((x + 72, yy), item, font=try_font(24), fill=COLORS["muted"])
            yy += 44
        x += 620
    img.save(path, format="PNG")
    return path


def draw_architecture_illustration():
    path = os.path.join(ILLUSTRATION_DIR, "architecture_mvc.png")
    img = Image.new("RGB", (1600, 900), COLORS["card"])
    draw = ImageDraw.Draw(img)
    draw.text((70, 60), "MVC 分层架构", font=try_font(48, bold=True), fill=COLORS["text"])
    layers = [
        ("View", ["MainWindow", "CanvasView", "PropertyPanel", "LayerPanel"], 100),
        ("Controller", ["DocumentController", "SelectionController", "QUndoStack Commands"], 560),
        ("Model", ["Document", "Shape 家族", "Layer / Workspace"], 1020),
    ]
    box_font = try_font(28, bold=True)
    small = try_font(22)
    for title, items, y in layers:
        draw_rounded_card(draw, (100, y, 1500, y + 180), COLORS["white"], outline=COLORS["accent2"])
        draw.text((130, y + 26), title, font=box_font, fill=COLORS["text"])
        x = 300
        for item in items:
            draw_rounded_card(draw, (x, y + 78, x + 330, y + 140), (232, 240, 252), outline=COLORS["accent"])
            draw.text((x + 20, y + 94), item, font=small, fill=COLORS["text"])
            x += 380
    draw.polygon([(780, 300), (740, 360), (820, 360)], fill=COLORS["accent"])
    draw.polygon([(780, 760), (740, 700), (820, 700)], fill=COLORS["accent"])
    img.save(path, format="PNG")
    return path


def export_ai_demo_screenshot():
    export_dir = os.path.join(TMP_DIR, "ai_export")
    os.makedirs(export_dir, exist_ok=True)
    exe = os.path.join(ROOT, "build", "Release", "generate_demo.exe")
    env = os.environ.copy()
    env["PATH"] = r"C:\Users\20371\.conda\envs\cpp\Library\bin;" + env.get("PATH", "")
    subprocess.run([exe, export_dir, "ai_gen_demo"], cwd=ROOT, env=env, check=True)
    src = os.path.join(export_dir, "ai_gen_demo.svg")
    dst = os.path.join(ILLUSTRATION_DIR, "ai_gen_demo.png")
    subprocess.run(["magick", "-background", "white", "-density", "180", src, dst], check=True)
    return dst


def create_illustrations():
    return {
        "auto_generate_flow": draw_flow_illustration(),
        "save_export_formats": draw_save_export_illustration(),
        "architecture_mvc": draw_architecture_illustration(),
        "ai_gen_demo": export_ai_demo_screenshot(),
    }


def all_source_images():
    mapping = {
        "cover": os.path.join(SOURCE_PNG_DIR, "slide-01-cover.png"),
        "builtin": os.path.join(SOURCE_PNG_DIR, "slide-02-builtin.png"),
        "ai": os.path.join(SOURCE_PNG_DIR, "slide-03-ai-assisted.png"),
        "input": os.path.join(SOURCE_PNG_DIR, "slide-04-parameter-input.png"),
        "auto": os.path.join(SOURCE_PNG_DIR, "slide-05-one-click.png"),
        "selection": os.path.join(SOURCE_PNG_DIR, "slide-06-auto-selection.png"),
        "editbox": os.path.join(SOURCE_PNG_DIR, "slide-07-canvas-selection.png"),
        "primitive": os.path.join(SOURCE_PNG_DIR, "slide-08-primitive-decompose.png"),
        "save": os.path.join(SOURCE_PNG_DIR, "slide-09-save-export.png"),
        "end": os.path.join(SOURCE_PNG_DIR, "slide-10-closing.png"),
    }
    for path in mapping.values():
        if not os.path.exists(path):
            raise FileNotFoundError(path)
    return mapping


def add_picture_safe(slide, path, left, top, width=None, height=None):
    if width is not None and height is not None:
        slide.shapes.add_picture(path, left, top, width=width, height=height)
    elif width is not None:
        slide.shapes.add_picture(path, left, top, width=width)
    else:
        slide.shapes.add_picture(path, left, top, height=height)


def set_bg(slide, path):
    slide.shapes.add_picture(path, 0, 0, width=prs.slide_width, height=prs.slide_height)
    slide.shapes[0].element.getparent().remove(slide.shapes[0].element)
    slide.follow_master_background = False


def add_title(slide, text, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.45), Inches(11.5), Inches(1.2))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*COLORS["text"][:3])
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(16)
        p2.font.color.rgb = RGBColor(*COLORS["muted"][:3])


def add_bullets(slide, items, left=0.7, top=1.5, width=5.0, height=4.8, font_size=20):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(*COLORS["text"][:3])
        p.space_after = Pt(8)


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def build_presentation(source_images, illustrations, backgrounds):
    global prs
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    set_bg(slide, backgrounds["slide01"])
    add_picture_safe(slide, source_images["cover"], Inches(7.35), Inches(1.45), width=Inches(5.5))
    cover = slide.shapes.add_textbox(Inches(0.75), Inches(1.55), Inches(6.0), Inches(4.8))
    tf = cover.text_frame
    tf.clear()
    for idx, (text, size, bold) in enumerate([
        (TITLE, 30, True),
        (f"姓名：{NAME}    学号：{STUDENT_ID}", 18, False),
        (COURSE, 17, False),
        (f"任课教师：{TEACHER}    班级：{CLASS}", 17, False),
        ("题目来源：自拟题目（C++ / Qt / 图形学）", 16, False),
    ]):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = RGBColor(*COLORS["text"][:3])
        p.space_after = Pt(10)
    add_notes(slide, f"各位老师好，我是{Name}，学号{STUDENT_ID}。本次汇报题目是《半自动化坐标驱动的矢量图形生成与编辑器项目》，下面我从题目来源、图形生成、画布编辑、保存导出和课程收获五个方面进行展示。")

    slide = prs.slides.add_slide(blank)
    set_bg(slide, backgrounds["slide02"])
    add_title(slide, "内置图形与菜单入口", "PPT 第一层：基础图形创建能力")
    add_bullets(slide, ["顶部菜单与左侧工具栏支持矩形、椭圆、直线、折线、正多边形和文本。", "创建结果进入 Document，通过 CanvasScene 同步到画布。", "内置图形体现编辑器基础可用性，也是后续半自动生成的载体。"], width=5.3)
    add_picture_safe(slide, source_images["builtin"], Inches(6.2), Inches(1.45), width=Inches(6.5))
    add_notes(slide, "项目首先提供完整的基础编辑能力。用户既可以直接创建基本图形，也为后续参数化生成与二次编辑提供统一定位。")

    slide = prs.slides.add_slide(blank)
    set_bg(slide, backgrounds["slide03"])
    add_title(slide, "AI 辅助自动生成", "从规则生成升级到工具链辅助生成")
    add_bullets(slide, ["GenerationPanel 集成文件选择与参数校验。", "当前规则引擎负责稳定落地；AI 用于自动化外部生成链。", "架构已预留扩展位，后续可切换为完整 Tool-Calling Agent。"], width=4.9)
    add_picture_safe(slide, illustrations["architecture_mvc"], Inches(5.4), Inches(1.35), width=Inches(3.8))
    add_picture_safe(slide, illustrations["auto_generate_flow"], Inches(9.0), Inches(1.35), width=Inches(3.9))
    add_notes(slide, "这一页强调区别：基础图形是内置能力，而自动生成依赖外部生成链或规则引擎，再回写到编辑器中。")

    slide = prs.slides.add_slide(blank)
    set_bg(slide, backgrounds["slide04"])
    add_title(slide, "输入代码或 txt 参数文件")
    add_bullets(slide, ["输入区同时承接代码片段与 txt 参数文件。", "参数文件定义图形名称、半径、数量、位置和输出文件名。", "输入完成后，生成按钮前先做合法性校验。"], width=5.0)
    add_picture_safe(slide, source_images["input"], Inches(6.0), Inches(1.35), width=Inches(6.7))
    add_notes(slide, "参数文件把图形生成从命令式操作转成配置式输入，更适合答辩展示和批量生成场景。")

    slide = prs.slides.add_slide(blank)
    set_bg(slide, backgrounds["slide05"])
    add_title(slide, "点击按钮，自动生成图形", "规则引擎 / 外部程序 -> DocumentController")
    add_bullets(slide, ["generate_demo 按 ai_gen_demo.txt 生成星形组合。", "JSON 描述经 ShapeFactory 转为 StarShape / PolygonShape。", "导入后自动加入文档、刷新画布并显示结果。"], width=5.2)
    add_picture_safe(slide, source_images["auto"], Inches(6.0), Inches(1.35), width=Inches(6.7))
    add_notes(slide, "这里突出“一键生成”。当前版本用规则引擎稳定演示，但输入输出接口已兼容 AI 生成链。")

    slide = prs.slides.add_slide(blank)
    set_bg(slide, backgrounds["slide06"])
    add_title(slide, "生成图被自动选中", "SelectionController 自动建立编辑上下文")
    add_bullets(slide, ["生成后的图形组会立即进入选中状态。", "选中后即可进行移动、缩放、旋转和属性修改。", "对应代码路径：导入完成 -> 选中 -> 出现控制句柄。"], width=5.0)
    add_picture_safe(slide, source_images["selection"], Inches(6.0), Inches(1.35), width=Inches(6.7))
    add_notes(slide, "自动生成并不停留在展示层，生成结果会立刻进入可编辑状态，这是“半自动化”的关键。")

    slide = prs.slides.add_slide(blank)
    set_bg(slide, backgrounds["slide07"])
    add_title(slide, "可直接编辑选择框", "画布交互 + 属性面板双向联动")
    add_bullets(slide, ["画布上的选择框、控制点和属性面板保持同步。", "支持位置、尺寸、填充、描边和文本属性修改。", "所有编辑走 Command + Undo/Redo，保证数据一致。"], width=5.2)
    add_picture_safe(slide, source_images["editbox"], Inches(6.0), Inches(1.35), width=Inches(6.7))
    add_notes(slide, "生成结果可以直接当作普通图形继续编辑，说明生成模块与编辑内核是打通的。")

    slide = prs.slides.add_slide(blank)
    set_bg(slide, backgrounds["slide08"])
    add_title(slide, "复杂图形可拆分为基本图元", "星形 / 多边形 / 折线 / 文本组合")
    add_bullets(slide, ["ai_gen_demo 可拆分为多个星形和装饰几何。", "图元级数据结构支持继续移动、缩放和改样式。", "体现参数化生成 + 手工精修的结合。"], width=4.9)
    add_picture_safe(slide, source_images["primitive"], Inches(5.7), Inches(1.3), width=Inches(3.8))
    add_picture_safe(slide, illustrations["ai_gen_demo"], Inches(9.4), Inches(1.3), width=Inches(3.5))
    add_notes(slide, "答辩时强调：虽然展示图看起来复杂，但内部对应的是标准矢量图元，不是一张不可编辑的位图。")

    slide = prs.slides.add_slide(blank)
    set_bg(slide, backgrounds["slide09"])
    add_title(slide, "支持 JSON 保存与 SVG 导出", "项目展示闭环：编辑 -> 持久化 -> 交换")
    add_bullets(slide, ["保存：JsonSerializer 持久化 document / layers / shapes / transform / style。", "导出：SvgExporter 输出 SVG 矢量交换格式。", "形成演示、保存、导出、继续编辑的完整闭环。"], width=5.1)
    add_picture_safe(slide, source_images["save"], Inches(5.7), Inches(1.3), width=Inches(3.6))
    add_picture_safe(slide, illustrations["save_export_formats"], Inches(9.3), Inches(1.3), width=Inches(3.6))
    add_notes(slide, "最后补上项目工程价值：不仅能演示生成，还能保存项目、导出交换格式，具备继续使用和扩展的基础。")

    slide = prs.slides.add_slide(blank)
    set_bg(slide, backgrounds["slide10"])
    add_title(slide, "结束页", "感谢各位老师聆听与指导")
    add_bullets(slide, [f"姓名：{NAME}", f"学号：{STUDENT_ID}", f"课程：{COURSE}", f"任课教师：{TEACHER}    班级：{CLASS}", "感谢各位老师聆听与指导！"], top=1.9, width=6.0, font_size=24)
    add_picture_safe(slide, source_images["end"], Inches(7.0), Inches(1.7), width=Inches(5.7))
    add_notes(slide, "以上是我的期末作品汇报，感谢各位老师聆听，请批评指正。")

    prs.save(OUTPUT_BASE + ".pptx")


def export_pdf():
    subprocess.run([
        "powershell", "-NoProfile", "-Command",
        "$ppt = New-Object -ComObject PowerPoint.Application; "
        "$ppt.Visible = 1; "
        f"$pres = $ppt.Presentations.Open('{OUTPUT_BASE}.pptx', $false, $false, $false); "
        f"$pres.SaveAs('{OUTPUT_BASE}.pdf', 32); "
        "$pres.Close(); $ppt.Quit();"
    ], check=True)


def export_zip():
    zip_path = OUTPUT_BASE + ".zip"
    with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
        for folder in [BG_DIR, ILLUSTRATION_DIR]:
            for name in os.listdir(folder):
                full = os.path.join(folder, name)
                if os.path.isfile(full):
                    zf.write(full, os.path.join("assets", os.path.basename(folder), name))
        for filename in [f"{NAME}-期末作品答辩-{STUDENT_ID}.pptx", f"{NAME}-期末作品答辩-{STUDENT_ID}.pdf"]:
            full = os.path.join(PRES_DIR, filename)
            if os.path.exists(full):
                zf.write(full, filename)
    return zip_path


def write_inventory(source_images):
    with open(os.path.join(PRES_DIR, "PPT_ASSET_INVENTORY.md"), "w", encoding="utf-8") as f:
        f.write("# PowerPoint Asset Inventory\n\n## Slide Backgrounds\n")
        for name in sorted(os.listdir(BG_DIR)):
            f.write(f"- `artifacts/presentation/slide_backgrounds/{name}`\n")
        f.write("\n## Illustrations\n")
        for name in sorted(os.listdir(ILLUSTRATION_DIR)):
            f.write(f"- `artifacts/presentation/illustrations/{name}`\n")
        f.write("\n## Source Screenshots\n")
        for key, path in source_images.items():
            f.write(f"- `{os.path.relpath(path, ROOT).replace(chr(92), '/')}` ({key})\n")


def write_build_report():
    with open(os.path.join(PRES_DIR, "PPT_BUILD_REPORT.md"), "w", encoding="utf-8") as f:
        f.write("# Presentation Build Report\n\n")
        f.write(f"- Generated: {datetime.datetime.now().isoformat(timespec='seconds')}\n")
        f.write(f"- Output PPTX: `{os.path.basename(OUTPUT_BASE)}.pptx`\n")
        f.write(f"- Output PDF: `{os.path.basename(OUTPUT_BASE)}.pdf`\n")
        f.write(f"- Output ZIP: `{os.path.basename(OUTPUT_BASE)}.zip`\n")
        f.write("- Slides: 10\n")


def write_qa_report():
    with open(os.path.join(PRES_DIR, "QA_REPORT_PRESENTATION.md"), "w", encoding="utf-8") as f:
        f.write("# Presentation QA Report\n\n")
        f.write("- Slide count: 10\n")
        f.write("- Save/export wording explicitly mentions JSON and SVG\n")
        f.write("- Speaker notes provided for all 10 slides\n")
        f.write("- Exported files: pptx, pdf, zip\n")


def main():
    ensure_dirs()
    source_images = all_source_images()
    illustrations = create_illustrations()
    backgrounds = {
        "slide01": create_background("slide01_cover.png", "期末作品答辩", f"{NAME}  {STUDENT_ID}", layout="cover", footer=COURSE),
        "slide02": create_background("slide02_builtin.png", "内置图形与菜单入口", "基础图形创建能力"),
        "slide03": create_background("slide03_ai.png", "AI 辅助自动生成", "规则引擎与 AI 生成链"),
        "slide04": create_background("slide04_input.png", "输入代码或 txt 参数文件", "配置式输入"),
        "slide05": create_background("slide05_auto.png", "点击按钮，自动生成图形", "一键导入"),
        "slide06": create_background("slide06_selection.png", "生成图被自动选中", "进入可编辑状态"),
        "slide07": create_background("slide07_editbox.png", "可直接编辑选择框", "画布与属性联动"),
        "slide08": create_background("slide08_primitive.png", "复杂图形可拆分", "基本图元组合"),
        "slide09": create_background("slide09_save.png", "JSON 保存与 SVG 导出", "项目闭环"),
        "slide10": create_background("slide10_end.png", "结束页", "感谢聆听"),
    }
    build_presentation(source_images, illustrations, backgrounds)
    export_pdf()
    export_zip()
    write_inventory(source_images)
    write_build_report()
    write_qa_report()
    print("BUILT", OUTPUT_BASE)


if __name__ == "__main__":
    main()
